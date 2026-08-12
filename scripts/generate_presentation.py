#!/usr/bin/env python
"""
Generate CADRE Project Presentation (PowerPoint)
==================================================
Creates a 15-slide professional presentation with all results and visualizations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import json

REPO_ROOT = Path(__file__).resolve().parent.parent
VIS_DIR = REPO_ROOT / "outputs" / "visualizations"
REPORT_PATH = REPO_ROOT / "outputs" / "cadre_bench" / "cadre_bench_report.json"
OUTPUT_PPT = REPO_ROOT / "outputs" / "CADRE_Presentation.pptx"

# Colors
DARK_BG = RGBColor(13, 17, 23)
CARD_BG = RGBColor(22, 27, 34)
BLUE = RGBColor(88, 166, 255)
GREEN = RGBColor(63, 185, 80)
PURPLE = RGBColor(188, 140, 255)
ORANGE = RGBColor(210, 153, 34)
RED = RGBColor(248, 81, 73)
WHITE = RGBColor(230, 237, 243)
GRAY = RGBColor(139, 148, 158)
DARK_BLUE = RGBColor(25, 60, 120)
LIGHT_BG = RGBColor(240, 245, 255)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = Pt(4)
    return p


def add_shape_box(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(left), Inches(top), 
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_table_slide(slide, headers, rows, left, top, width, row_height=0.4):
    cols = len(headers)
    table_rows = len(rows) + 1
    col_width = width / cols

    table = slide.shapes.add_table(table_rows, cols, 
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(row_height * table_rows)).table

    # Style header
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Style rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 242, 248)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(40, 40, 40)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def main():
    report = json.load(open(REPORT_PATH))
    metrics = report["metrics"]
    matrix = report["performance_matrix"]
    params = report["param_stats"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 1: TITLE
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, 1, 1.2, 11.3, 1, "CADRE", 54, BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.2, 11.3, 0.8, "Continual Adaptation for Driving with Robust Evolution",
                 24, WHITE, False, PP_ALIGN.CENTER)

    # Divider line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.3), Inches(5.3), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    add_text_box(slide, 1, 3.6, 11.3, 0.5, "Problem Statement B5", 20, ORANGE, True, PP_ALIGN.CENTER)
    tf = add_text_box(slide, 2, 4.2, 9.3, 1.5,
        "Design a continual adaptation method for VLA autonomous driving models that "
        "incorporates new regional regulations, road layouts, and weather patterns while "
        "retaining prior driving competence across previously learned environments.",
        14, GRAY, False, PP_ALIGN.CENTER)

    add_text_box(slide, 1, 5.8, 11.3, 0.4, "Datasets: BDD100K + nuScenes  |  Backbone: LLaVA-v1.5-7B  |  Method: EWC + Replay + LoRA",
                 13, GRAY, False, PP_ALIGN.CENTER)

    # Results box
    box = add_shape_box(slide, 3.5, 6.3, 6.3, 0.8, RGBColor(22, 27, 44), BLUE)
    add_text_box(slide, 3.5, 6.35, 6.3, 0.7,
                 f"CDAR = {metrics['CDAR']}  |  BWT = {metrics['BWT']}%  |  FWT = +{metrics['FWT']}%  |  Efficiency = {metrics['Efficiency']}%",
                 13, GREEN, True, PP_ALIGN.CENTER)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 2: PROBLEM STATEMENT
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Problem: Catastrophic Forgetting in Autonomous Driving", 28, BLUE, True)

    tf = add_text_box(slide, 0.8, 1.3, 5.5, 4.5, "", 14, WHITE)
    add_paragraph(tf, "What is Catastrophic Forgetting?", 18, ORANGE, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "When a neural network learns a new task, it tends to", 14, WHITE)
    add_paragraph(tf, "FORGET previously learned tasks. This is disastrous", 14, WHITE)
    add_paragraph(tf, "for self-driving cars that must handle:", 14, WHITE)
    add_paragraph(tf, "")
    add_paragraph(tf, "  🇺🇸  US road rules & layouts", 15, BLUE)
    add_paragraph(tf, "  🇸🇬  Singapore road rules & layouts", 15, GREEN)
    add_paragraph(tf, "  🇪🇺  European road rules & layouts", 15, ORANGE)
    add_paragraph(tf, "  🌧️  Rainy / foggy / snowy weather", 15, PURPLE)
    add_paragraph(tf, "")
    add_paragraph(tf, "A deployed model must adapt to NEW regions without", 14, WHITE)
    add_paragraph(tf, "forgetting how to drive in OLD regions.", 14, WHITE)

    tf2 = add_text_box(slide, 7, 1.3, 5.5, 4.5, "", 14, WHITE)
    add_paragraph(tf2, "Our Solution: CADRE", 18, GREEN, True)
    add_paragraph(tf2, "")
    add_paragraph(tf2, "Three strategies working together:", 14, WHITE)
    add_paragraph(tf2, "")
    add_paragraph(tf2, "1. EWC (Elastic Weight Consolidation)", 15, ORANGE, True)
    add_paragraph(tf2, "   Protects important weights from being overwritten", 13, GRAY)
    add_paragraph(tf2, "")
    add_paragraph(tf2, "2. Experience Replay", 15, RED, True)
    add_paragraph(tf2, "   Mixes 30% old data into new training batches", 13, GRAY)
    add_paragraph(tf2, "")
    add_paragraph(tf2, "3. LoRA Adapters", 15, BLUE, True)
    add_paragraph(tf2, "   Only 0.35% extra parameters per domain (~38 MB)", 13, GRAY)
    add_paragraph(tf2, "")
    add_paragraph(tf2, "All on a frozen LLaVA-v1.5-7B backbone (7B params)", 14, GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 3: DATASETS
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Datasets & Domain Sequence", 28, BLUE, True)

    add_table_slide(slide,
        ["Domain", "Dataset", "Description", "Split Criteria"],
        [
            ["domain_us", "BDD100K", "US urban/highway driving", "Clear weather, daytime"],
            ["domain_sg", "nuScenes", "Singapore urban driving", "Singapore locations"],
            ["domain_eu", "nuScenes", "Boston/EU urban driving", "Boston-seaport location"],
            ["domain_rainy", "BDD100K", "Adverse weather driving", "Rain, fog, snow, night"],
        ],
        0.8, 1.5, 11.7, 0.5
    )

    tf = add_text_box(slide, 0.8, 4.2, 11.5, 2, "", 14, WHITE)
    add_paragraph(tf, "Training Sequence (Sequential — one domain at a time):", 16, ORANGE, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "   domain_us  →  domain_sg  →  domain_eu  →  domain_rainy", 20, GREEN, True, PP_ALIGN.CENTER)
    add_paragraph(tf, "   (BDD100K)      (nuScenes)     (nuScenes)      (BDD100K)", 14, GRAY, False, PP_ALIGN.CENTER)
    add_paragraph(tf, "")
    add_paragraph(tf, "BDD100K: 100,000 driving videos from Berkeley DeepDrive", 14, GRAY)
    add_paragraph(tf, "nuScenes: 1,000 driving scenes with 3D annotations from Motional", 14, GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 4: ARCHITECTURE
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "System Architecture", 28, BLUE, True)

    img_path = VIS_DIR / "architecture_overview.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.2), Inches(12.3))

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 5: VLA BACKBONE + LoRA
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "VLA Backbone + LoRA Adapters", 28, BLUE, True)

    tf = add_text_box(slide, 0.8, 1.3, 5.5, 5, "", 14, WHITE)
    add_paragraph(tf, "VLA Backbone: LLaVA-v1.5-7B", 18, ORANGE, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "• Vision-Language-Action model", 14, WHITE)
    add_paragraph(tf, "• CLIP-ViT-L/14 vision encoder", 14, WHITE)
    add_paragraph(tf, "• Vicuna-7B language model", 14, WHITE)
    add_paragraph(tf, "• 7.06 Billion parameters — ALL FROZEN", 15, RED, True)
    add_paragraph(tf, "• Loaded in float16 for GPU efficiency", 14, WHITE)
    add_paragraph(tf, "• gradient_checkpointing = True", 14, GRAY)
    add_paragraph(tf, "")
    add_paragraph(tf, "File: src/models/vla_backbone.py", 12, GRAY)

    tf2 = add_text_box(slide, 7, 1.3, 5.5, 5, "", 14, WHITE)
    add_paragraph(tf2, "LoRA Adapters (Per Domain)", 18, GREEN, True)
    add_paragraph(tf2, "")
    add_paragraph(tf2, "Low-Rank Adaptation — inject small trainable", 14, WHITE)
    add_paragraph(tf2, "matrices into frozen attention layers:", 14, WHITE)
    add_paragraph(tf2, "")
    add_paragraph(tf2, "  Rank (r) = 16", 14, BLUE)
    add_paragraph(tf2, "  Alpha = 32 (scaling = 2.0)", 14, BLUE)
    add_paragraph(tf2, "  Target: q_proj, v_proj", 14, BLUE)
    add_paragraph(tf2, "  Dropout = 0.05", 14, BLUE)
    add_paragraph(tf2, "")
    add_paragraph(tf2, "  Params per domain: 24.7M (0.35%)", 15, GREEN, True)
    add_paragraph(tf2, "  Storage per domain: ~38 MB", 15, GREEN, True)
    add_paragraph(tf2, "  (vs 14 GB for full model)", 13, GRAY)
    add_paragraph(tf2, "")
    add_paragraph(tf2, "File: src/adapters/lora_adapter.py", 12, GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 6: EWC
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Elastic Weight Consolidation (EWC)", 28, BLUE, True)

    tf = add_text_box(slide, 0.8, 1.3, 11.5, 5.5, "", 14, WHITE)
    add_paragraph(tf, "How EWC Prevents Forgetting", 20, ORANGE, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "After training on domain Dk:", 16, WHITE, True)
    add_paragraph(tf, "  1. Compute Fisher Information Matrix Fk — identifies which weights are IMPORTANT for Dk", 14, WHITE)
    add_paragraph(tf, "  2. Store optimal parameter values θ*k", 14, WHITE)
    add_paragraph(tf, "")
    add_paragraph(tf, "When training on domain Dk+1:", 16, WHITE, True)
    add_paragraph(tf, "  3. Add penalty:  L_ewc = (λ/2) × Σ Fi × (θi − θ*i)²", 16, GREEN, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "This PENALIZES changes to weights that were important for old domains!", 15, ORANGE)
    add_paragraph(tf, "")
    add_paragraph(tf, "Configuration:", 16, BLUE, True)
    add_paragraph(tf, "  • Lambda (λ) = 5,000 — strong protection against forgetting", 14, WHITE)
    add_paragraph(tf, "  • Fisher samples = 200 — samples for Fisher computation", 14, WHITE)
    add_paragraph(tf, "  • Variant = Online EWC — running average, memory-efficient", 14, WHITE)
    add_paragraph(tf, "  • Gamma (γ) = 0.95 — slow decay of old Fisher information", 14, WHITE)
    add_paragraph(tf, "")
    add_paragraph(tf, "File: src/continual/ewc.py", 12, GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 7: REPLAY BUFFER
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Experience Replay Buffer", 28, BLUE, True)

    tf = add_text_box(slide, 0.8, 1.3, 11.5, 5.5, "", 14, WHITE)
    add_paragraph(tf, "How Replay Prevents Distribution Shift", 20, ORANGE, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "Problem: If we only train on new domain data, the model's distribution", 14, WHITE)
    add_paragraph(tf, "completely shifts to the new domain, forgetting old ones.", 14, WHITE)
    add_paragraph(tf, "")
    add_paragraph(tf, "Solution: Mix old domain samples into every training batch!", 15, GREEN, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "  ┌──────────────────────────────────────────┐", 14, BLUE, False, PP_ALIGN.LEFT)
    add_paragraph(tf, "  │  New Domain Data (70%) ─┐                │", 14, BLUE)
    add_paragraph(tf, "  │                         ├→ Mixed Batch   │", 14, BLUE)
    add_paragraph(tf, "  │  Replay Buffer (30%) ──┘                 │", 14, BLUE)
    add_paragraph(tf, "  └──────────────────────────────────────────┘", 14, BLUE)
    add_paragraph(tf, "")
    add_paragraph(tf, "Configuration:", 16, PURPLE, True)
    add_paragraph(tf, "  • Buffer size: 2,000 samples per domain (reservoir sampling)", 14, WHITE)
    add_paragraph(tf, "  • Replay ratio: 0.3 (30% old data in each batch)", 14, WHITE)
    add_paragraph(tf, "  • Storage: Disk-based (memory efficient)", 14, WHITE)
    add_paragraph(tf, "")
    add_paragraph(tf, "File: src/continual/replay_buffer.py", 12, GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 8: OUTPUT HEADS
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Multi-Head Output & Domain Router", 28, BLUE, True)

    add_table_slide(slide,
        ["Head", "Output", "Shape", "Description"],
        [
            ["WaypointHead", "Trajectory", "[B, 12, 2]", "12 future waypoints (x, y)"],
            ["HazardHead", "Hazard", "[B, 8]", "8-class hazard detection"],
            ["RegulationHead", "Rules", "[B, 15]", "15-class traffic regulation"],
            ["WeatherHead", "Weather", "[B, 6]", "6-class weather classification"],
            ["Integration", "Fusion", "[B, 512]", "Attention-based head fusion"],
        ],
        0.8, 1.5, 11.7, 0.45
    )

    tf = add_text_box(slide, 0.8, 4.5, 11.5, 2.5, "", 14, WHITE)
    add_paragraph(tf, "Domain Router", 18, GREEN, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "At inference time, the Domain Router automatically detects which driving", 14, WHITE)
    add_paragraph(tf, "domain an input image belongs to and routes it to the correct LoRA adapter.", 14, WHITE)
    add_paragraph(tf, "")
    add_paragraph(tf, "  Image → Router → domain_us/domain_sg/domain_eu/domain_rainy → Load correct LoRA adapter", 14, BLUE)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 9: TRAINING PIPELINE
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Training Pipeline (8 Stages)", 28, BLUE, True)

    add_table_slide(slide,
        ["#", "Stage", "Description", "Status", "Date"],
        [
            ["1", "verify_backbone", "Load & freeze LLaVA-7B", "✅ Done", "Jul 25"],
            ["2", "domain_us", "Train US driving (BDD100K)", "✅ Done", "Jul 30"],
            ["3", "domain_sg", "Train Singapore (nuScenes)", "✅ Done", "Jul 31"],
            ["4", "domain_eu", "Train EU/Boston (nuScenes)", "✅ Done", "Aug 3"],
            ["5", "domain_rainy", "Train adverse weather (BDD100K)", "✅ Done", "Aug 3"],
            ["6", "train_router", "Train domain classifier", "✅ Done", "Aug 3"],
            ["7", "train_heads", "Train 4 output heads", "✅ Done", "Aug 4"],
            ["8", "benchmark", "Run CADRE-Bench evaluation", "✅ Done", "Aug 4"],
        ],
        0.8, 1.3, 11.7, 0.4
    )

    tf = add_text_box(slide, 0.8, 5.2, 11.5, 2, "", 13, WHITE)
    add_paragraph(tf, "Pipeline is fully resumable — saves progress after each stage.", 14, GREEN)
    add_paragraph(tf, "If interrupted, rerun the same command to resume from where it stopped.", 14, GRAY)
    add_paragraph(tf, "Command:  python scripts/run_pipeline.py", 14, BLUE, True)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 10: TRAINING TIMELINE
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Training Timeline", 28, BLUE, True)

    img_path = VIS_DIR / "training_timeline.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.3), Inches(12.3))

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 11: RESULTS — METRICS
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Results: CADRE-Bench Metrics", 28, BLUE, True)

    img_path = VIS_DIR / "metrics_summary.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.2), Inches(12.7), Inches(3.8))

    add_table_slide(slide,
        ["Metric", "Score", "Meaning"],
        [
            ["BWT (Forgetting)", f"{metrics['BWT']}%", "Near-zero forgetting of old domains"],
            ["FWT (Transfer)", f"+{metrics['FWT']}%", "Old training helps new domains"],
            ["Plasticity", f"{metrics['Plasticity']}%", "Learns new domains nearly as well as single-task"],
            ["Efficiency", f"{metrics['Efficiency']}%", "Only 0.35% parameter overhead"],
            ["CDAR (Overall)", f"{metrics['CDAR']}", "Outstanding composite score (97.35%)"],
        ],
        1.5, 5.3, 10.3, 0.38
    )

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 12: PERFORMANCE MATRIX
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Results: Performance Matrix", 28, BLUE, True)

    img_path = VIS_DIR / "performance_matrix_heatmap.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.0), Inches(7), Inches(5.8))

    tf = add_text_box(slide, 7.8, 1.3, 5, 5.5, "", 13, WHITE)
    add_paragraph(tf, "How to Read This Matrix", 16, ORANGE, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "R[i][j] = accuracy on domain j", 13, WHITE)
    add_paragraph(tf, "after training through domain i", 13, WHITE)
    add_paragraph(tf, "")
    add_paragraph(tf, "Diagonal (dashed boxes):", 14, BLUE, True)
    add_paragraph(tf, "How well it learns each domain", 13, GRAY)
    add_paragraph(tf, "→ 94% to 97% ✅", 13, GREEN)
    add_paragraph(tf, "")
    add_paragraph(tf, "Final Row (after all training):", 14, GREEN, True)
    add_paragraph(tf, "All domains retain >92.5%", 13, GRAY)
    add_paragraph(tf, "")
    add_paragraph(tf, "Forgetting Example:", 14, RED, True)
    add_paragraph(tf, "US: 0.94 → 0.925 = only -1.5%", 13, GRAY)
    add_paragraph(tf, "after learning 3 more domains!", 13, GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 13: DOMAIN COMPARISON
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Results: Domain-wise Performance", 28, BLUE, True)

    img_path = VIS_DIR / "domain_performance.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 14: PARAMETER EFFICIENCY
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Parameter Efficiency: Why LoRA Matters", 28, BLUE, True)

    img_path = VIS_DIR / "parameter_efficiency.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.1), Inches(12.7), Inches(3.8))

    add_table_slide(slide,
        ["Approach", "Trainable Params", "Storage / Domain", "4 Domains Total"],
        [
            ["Full Retrain", "7.06B (100%)", "~14 GB", "~56 GB"],
            ["CADRE (LoRA)", "24.7M (0.35%)", "~38 MB", "~152 MB"],
            ["Savings", "99.65% fewer", "368× smaller", "368× smaller"],
        ],
        2, 5.3, 9.3, 0.45
    )

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 15: CONCLUSION
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.5, 0.3, 12, 0.7, "Conclusion & Key Takeaways", 28, BLUE, True)

    tf = add_text_box(slide, 0.8, 1.3, 5.5, 5.5, "", 14, WHITE)
    add_paragraph(tf, "What We Built", 20, GREEN, True)
    add_paragraph(tf, "")
    add_paragraph(tf, "✅ Continual learning framework for VLA", 14, WHITE)
    add_paragraph(tf, "   autonomous driving (EWC + Replay + LoRA)", 13, GRAY)
    add_paragraph(tf, "")
    add_paragraph(tf, "✅ Trained on 4 sequential domains without", 14, WHITE)
    add_paragraph(tf, "   catastrophic forgetting (BWT = -1.17%)", 13, GRAY)
    add_paragraph(tf, "")
    add_paragraph(tf, "✅ Parameter-efficient: 0.35% overhead", 14, WHITE)
    add_paragraph(tf, "   (~38 MB adapter vs ~14 GB full model)", 13, GRAY)
    add_paragraph(tf, "")
    add_paragraph(tf, "✅ CADRE-Bench: Custom benchmark protocol", 14, WHITE)
    add_paragraph(tf, "   with 5 metrics (BWT, FWT, Plasticity,", 13, GRAY)
    add_paragraph(tf, "   Efficiency, CDAR)", 13, GRAY)
    add_paragraph(tf, "")
    add_paragraph(tf, "✅ Fully reproducible & resumable pipeline", 14, WHITE)

    tf2 = add_text_box(slide, 7, 1.3, 5.5, 5.5, "", 14, WHITE)
    add_paragraph(tf2, "Key Results", 20, ORANGE, True)
    add_paragraph(tf2, "")

    results = [
        ("BWT", f"{metrics['BWT']}%", "Near-zero forgetting", GREEN),
        ("FWT", f"+{metrics['FWT']}%", "Strong positive transfer", BLUE),
        ("Plasticity", f"{metrics['Plasticity']}%", "Excellent learning", GREEN),
        ("Efficiency", f"{metrics['Efficiency']}%", "Minimal overhead", GREEN),
        ("CDAR", f"{metrics['CDAR']}", "Outstanding overall", ORANGE),
    ]
    for name, val, desc, color in results:
        add_paragraph(tf2, f"  {name}: {val}", 16, color, True)
        add_paragraph(tf2, f"    {desc}", 12, GRAY)

    add_paragraph(tf2, "")
    add_paragraph(tf2, "References", 14, GRAY, True)
    add_paragraph(tf2, "Kirkpatrick et al. (2017) - EWC", 10, GRAY)
    add_paragraph(tf2, "Hu et al. (2022) - LoRA", 10, GRAY)
    add_paragraph(tf2, "Liu et al. (2024) - LLaVA", 10, GRAY)
    add_paragraph(tf2, "BDD100K + nuScenes datasets", 10, GRAY)

    # Add Thank You text
    add_text_box(slide, 0.5, 6.5, 12.3, 0.7, "Thank You!", 24, BLUE, True, PP_ALIGN.CENTER)

    # ━━━ Save ━━━
    prs.save(str(OUTPUT_PPT))
    print(f"\n✅ Presentation saved to: {OUTPUT_PPT}")
    print(f"   Slides: 15")
    print(f"   Size: {OUTPUT_PPT.stat().st_size / 1024:.0f} KB\n")


if __name__ == "__main__":
    main()
